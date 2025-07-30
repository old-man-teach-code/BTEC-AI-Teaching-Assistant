import logging
import os
from typing import List, Dict, Any, Optional
from pathlib import Path
import asyncio

# Document parsing libraries
import pypdf
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd

# LlamaIndex imports
from llama_index.core import Document as LlamaDocument
from llama_index.core.node_parser import SimpleNodeParser
from llama_index.core.text_splitter import SentenceSplitter

from core.config import settings
from database.vector_store import vector_store_manager
from utils.file_utils import get_file_extension, read_text_file

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """
    Service xử lý documents và extract knowledge
    Hỗ trợ các định dạng: PDF, DOCX, PPTX, TXT
    """
    
    def __init__(self):
        # Text splitter cho việc chia document thành chunks
        self.text_splitter = SentenceSplitter(
            chunk_size=settings.chunk_size,
            chunk_overlap=settings.chunk_overlap,
            separator=" ",
            paragraph_separator="\n\n",
            secondary_chunking_regex="[.!?]"  # Chia theo câu
        )
        
        # Node parser cho LlamaIndex
        self.node_parser = SimpleNodeParser.from_defaults(
            text_splitter=self.text_splitter,
            include_metadata=True,
            include_prev_next_rel=True  # Include relationships
        )
    
    async def process_document(
        self,
        file_path: str,
        user_id: int,
        metadata: Optional[Dict[str, Any]] = None
    ) -> int:
        """
        Process document và extract knowledge
        
        Args:
            file_path: Đường dẫn đến file document
            user_id: ID của user sở hữu document
            metadata: Metadata bổ sung
            
        Returns:
            Số chunks đã extract
            
        Raises:
            Exception: Nếu có lỗi trong quá trình xử lý
        """
        try:
            # Kiểm tra file tồn tại
            if not os.path.exists(file_path):
                raise FileNotFoundError(f"File không tồn tại: {file_path}")
            
            # Extract text từ document
            logger.info(f"Đang extract text từ: {file_path}")
            text_content = await self._extract_text(file_path)
            
            if not text_content or len(text_content.strip()) < 10:
                logger.warning(f"Document rỗng hoặc quá ngắn: {file_path}")
                return 0
            
            # Tạo LlamaIndex document
            doc_metadata = metadata or {}
            doc_metadata.update({
                "file_path": file_path,
                "file_name": os.path.basename(file_path),
                "user_id": str(user_id),
                "extraction_time": str(asyncio.get_event_loop().time())
            })
            
            llama_doc = LlamaDocument(
                text=text_content,
                metadata=doc_metadata
            )
            
            # Parse document thành nodes/chunks
            logger.info("Đang chia document thành chunks...")
            nodes = self.node_parser.get_nodes_from_documents([llama_doc])
            
            if not nodes:
                logger.warning("Không thể tạo chunks từ document")
                return 0
            
            logger.info(f"Đã tạo {len(nodes)} chunks")
            
            # Chuẩn bị data cho vector store
            documents_to_add = []
            for i, node in enumerate(nodes):
                # Thêm metadata cho mỗi chunk
                chunk_metadata = {
                    **node.metadata,
                    "chunk_index": i,
                    "total_chunks": len(nodes),
                    "chunk_id": f"{metadata.get('document_id', 'unknown')}_{i}"
                }
                
                # Nếu node có relationships (prev/next), thêm vào metadata
                if node.relationships:
                    chunk_metadata["has_relationships"] = True
                
                documents_to_add.append({
                    "text": node.get_content(),
                    "metadata": chunk_metadata,
                    "id": f"doc_{user_id}_{metadata.get('document_id', 'unknown')}_{i}"
                })
            
            # Thêm vào vector store
            logger.info("Đang thêm chunks vào vector store...")
            
            # Lưu vào cả user collection và global collection
            # User collection
            user_collection = vector_store_manager.get_user_collection_name(user_id)
            vector_store_manager.add_documents(
                documents=documents_to_add,
                collection_name=user_collection
            )
            
            # Global collection với user_id trong metadata
            vector_store_manager.add_documents(
                documents=documents_to_add,
                collection_name=settings.chroma_collection_name
            )
            
            logger.info(f"✅ Hoàn thành xử lý document: {len(documents_to_add)} chunks")
            return len(documents_to_add)
            
        except Exception as e:
            logger.error(f"❌ Lỗi khi process document {file_path}: {str(e)}")
            raise
    
    async def _extract_text(self, file_path: str) -> str:
        """
        Extract text từ document dựa vào file type
        
        Args:
            file_path: Đường dẫn file
            
        Returns:
            Text content của document
        """
        ext = get_file_extension(file_path).lower()
        
        try:
            if ext == '.pdf':
                return await self._extract_pdf(file_path)
            elif ext == '.docx':
                return await self._extract_docx(file_path)
            elif ext == '.pptx':
                return await self._extract_pptx(file_path)
            elif ext == '.txt':
                return await read_text_file(file_path)
            else:
                logger.warning(f"Định dạng file không được hỗ trợ: {ext}")
                return ""
                
        except Exception as e:
            logger.error(f"Lỗi extract text từ {file_path}: {str(e)}")
            raise
    
    async def _extract_pdf(self, file_path: str) -> str:
        """
        Extract text từ PDF file
        
        Args:
            file_path: Đường dẫn PDF file
            
        Returns:
            Extracted text
        """
        try:
            text_parts = []
            
            with open(file_path, 'rb') as file:
                pdf_reader = pypdf.PdfReader(file)
                
                # Extract từng page
                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            # Thêm page number để tracking
                            text_parts.append(f"[Trang {page_num + 1}]\n{page_text}")
                    except Exception as e:
                        logger.warning(f"Không thể extract trang {page_num + 1}: {str(e)}")
                        continue
            
            # Join all pages
            full_text = "\n\n".join(text_parts)
            
            # Clean up text
            full_text = self._clean_text(full_text)
            
            logger.info(f"Extracted {len(full_text)} characters từ PDF")
            return full_text
            
        except Exception as e:
            logger.error(f"Lỗi extract PDF: {str(e)}")
            raise
    
    async def _extract_docx(self, file_path: str) -> str:
        """
        Extract text từ DOCX file
        
        Args:
            file_path: Đường dẫn DOCX file
            
        Returns:
            Extracted text
        """
        try:
            doc = DocxDocument(file_path)
            text_parts = []
            
            # Extract từng paragraph
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            # Extract từ tables nếu có
            for table in doc.tables:
                table_text = self._extract_table_text(table)
                if table_text:
                    text_parts.append(f"\n[Bảng]\n{table_text}\n")
            
            # Join all parts
            full_text = "\n\n".join(text_parts)
            
            # Clean up text
            full_text = self._clean_text(full_text)
            
            logger.info(f"Extracted {len(full_text)} characters từ DOCX")
            return full_text
            
        except Exception as e:
            logger.error(f"Lỗi extract DOCX: {str(e)}")
            raise
    
    async def _extract_pptx(self, file_path: str) -> str:
        """
        Extract text từ PPTX file
        
        Args:
            file_path: Đường dẫn PPTX file
            
        Returns:
            Extracted text
        """
        try:
            prs = Presentation(file_path)
            text_parts = []
            
            # Extract từng slide
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                
                # Extract text từ shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                
                if slide_text:
                    slide_content = "\n".join(slide_text)
                    text_parts.append(f"[Slide {slide_num + 1}]\n{slide_content}")
            
            # Join all slides
            full_text = "\n\n".join(text_parts)
            
            # Clean up text
            full_text = self._clean_text(full_text)
            
            logger.info(f"Extracted {len(full_text)} characters từ PPTX")
            return full_text
            
        except Exception as e:
            logger.error(f"Lỗi extract PPTX: {str(e)}")
            raise
    
    def _extract_table_text(self, table) -> str:
        """
        Extract text từ table trong DOCX
        
        Args:
            table: Table object từ python-docx
            
        Returns:
            Formatted table text
        """
        try:
            rows_text = []
            
            for row in table.rows:
                cells_text = []
                for cell in row.cells:
                    cells_text.append(cell.text.strip())
                
                if any(cells_text):  # Chỉ add row nếu có content
                    rows_text.append(" | ".join(cells_text))
            
            return "\n".join(rows_text)
            
        except Exception as e:
            logger.warning(f"Lỗi extract table: {str(e)}")
            return ""
    
    def _clean_text(self, text: str) -> str:
        """
        Clean và normalize text
        
        Args:
            text: Raw text
            
        Returns:
            Cleaned text
        """
        # Remove multiple spaces
        text = " ".join(text.split())
        
        # Remove multiple newlines nhưng giữ paragraph breaks
        lines = text.split('\n')
        cleaned_lines = []
        
        for line in lines:
            line = line.strip()
            if line:
                cleaned_lines.append(line)
            elif cleaned_lines and cleaned_lines[-1] != '':
                # Add empty line for paragraph break
                cleaned_lines.append('')
        
        text = '\n'.join(cleaned_lines)
        
        # Remove các ký tự đặc biệt không cần thiết
        # Giữ lại các ký tự tiếng Việt
        text = text.replace('\x00', '')  # Null characters
        text = text.replace('\u200b', '')  # Zero-width space
        
        return text
    
    async def process_batch(
        self,
        file_paths: List[str],
        user_id: int,
        metadata_list: Optional[List[Dict[str, Any]]] = None
    ) -> Dict[str, Any]:
        """
        Process nhiều documents cùng lúc
        
        Args:
            file_paths: Danh sách file paths
            user_id: ID của user
            metadata_list: Danh sách metadata tương ứng
            
        Returns:
            Batch processing results
        """
        results = {
            "total": len(file_paths),
            "successful": 0,
            "failed": 0,
            "details": []
        }
        
        # Process từng document
        for i, file_path in enumerate(file_paths):
            metadata = metadata_list[i] if metadata_list and i < len(metadata_list) else {}
            
            try:
                chunks_count = await self.process_document(
                    file_path=file_path,
                    user_id=user_id,
                    metadata=metadata
                )
                
                results["successful"] += 1
                results["details"].append({
                    "file_path": file_path,
                    "success": True,
                    "chunks_extracted": chunks_count
                })
                
            except Exception as e:
                logger.error(f"Lỗi process {file_path}: {str(e)}")
                results["failed"] += 1
                results["details"].append({
                    "file_path": file_path,
                    "success": False,
                    "error": str(e)
                })
        
        return results